from io import BytesIO
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu

from src.api.v1.documents import _embed_pptx_preview_assets
from src.formats.pptx.agents.nodes.html_generator import (
    _build_human_content_with_images,
    _build_user_reference_images,
    _element_placements_for_html_prompt,
    _inject_fixed_template,
    _inject_visual_asset_images,
    _materialize_html_icon_node,
    _user_reference_images_for_slide,
)
from src.formats.pptx.agents.nodes.render_convert import (
    PPTX_TEXT_ALIGNMENT_PREVIEW_CSS,
    _embed_cached_icons,
    _embed_local_images,
    _normalize_legacy_icon_nodes,
    _normalize_slide_html,
)
from src.formats.pptx.agents.nodes.unified_planner import (
    _apply_layout_inset,
    _compact_parent_slide_plan,
    _extend_blueprints_to_requested_count,
    _extract_requested_slide_count,
    _extract_slide_revision_instructions,
    _infer_target_slide_count,
    _invoke_planner_with_structured_output,
    _normalize_blueprints,
    _normalize_revision_scope,
    _repair_missing_slide_count_result,
    _structured_planner_llm,
)
from src.formats.pptx.agents.nodes.visual_asset_planner import (
    DIAGRAMS_DPI,
    METHOD_DIAGRAMS,
    MIN_DIAGRAM_PNG_HEIGHT,
    MIN_DIAGRAM_PNG_WIDTH,
    _add_diagram_image_border,
    _diagrams_node_candidates,
    _diagrams_topology,
    _fallback_diagrams_edges,
    _fallback_diagrams_nodes,
    _fallback_plan,
    _fallback_visual_slot_plan,
    _fit_png_canvas_to_placement,
    _layout_nodes,
    _llm_plan_validation_issues,
    _looks_like_truncated_json,
    _merge_missing_reserved_slot_assets,
    _negative_visual_asset_signal,
    _normalize_plan,
    _parse_mermaid,
    _quick_visual_signal,
    _render_asset,
    _safe_diagrams_topology,
    _slide_context_requires_diagram,
    _slot_constraints,
)
from src.formats.pptx.mapper.engine import CSStoOOXMLEngine
from src.formats.pptx.mapper.html_parser import ParsedElement, parse_slide_html
from src.formats.pptx.rulesets import get_ruleset
from src.formats.pptx.rulesets.validator import DesignQualityEvaluator
from src.utils.iconify import get_fallback_icon_path
from src.utils.image_gen import _generate_with_nova_canvas
from src.utils.json_repair import parse_llm_json


def _placements_overlap_for_test(placements: list[dict]) -> bool:
    boxed = [
        item for item in placements
        if all(key in item for key in ("x", "y", "w", "h"))
        and str(item.get("element", "")).lower() not in {"line", "connector"}
    ]
    for index, first in enumerate(boxed):
        for second in boxed[index + 1:]:
            ax, ay, aw, ah = first["x"], first["y"], first["w"], first["h"]
            bx, by, bw, bh = second["x"], second["y"], second["w"], second["h"]
            x_overlap = max(0, min(ax + aw, bx + bw) - max(ax, bx))
            y_overlap = max(0, min(ay + ah, by + bh) - max(ay, by))
            if x_overlap <= 0 or y_overlap <= 0:
                continue
            ratio = (x_overlap * y_overlap) / max(1, min(aw * ah, bw * bh))
            if ratio > 0.08:
                return True
    return False


def test_standard_layout_catalog_exposes_master_zones_and_body_patterns() -> None:
    ruleset = get_ruleset()

    assert len(ruleset.layout_zones["header_zones"]) >= 20
    assert len(ruleset.layout_zones["footer_zones"]) >= 20
    assert len(ruleset.icon_layouts["placements"]) >= 20
    assert "data-pptx-type=\"icon\"" in ruleset.get_icon_layout_rules()
    assert (
        ruleset.layout_patterns["total_patterns"] >= 100
        and ruleset.layout_patterns["total_patterns"] <= 300
    )
    assert sum(
        len(category["patterns"])
        for category in ruleset.layout_patterns["categories"].values()
    ) == ruleset.layout_patterns["total_patterns"]
    assert ruleset.get_body_layout("compare_vs__evidence")["variant_of"] == "compare_vs"
    rules = ruleset.get_planner_layout_rules()
    assert "OOXML Planning Boundary" in rules
    assert "header_claim_bar" in rules
    assert "dashboard_chart_sidebar" in rules


def test_revision_request_extracts_per_slide_instructions() -> None:
    instructions = _extract_slide_revision_instructions(
        "슬라이드 3: LangChain RAG Pipeline 아키텍처 다이어그램으로 변경해줘. "
        "슬라이드 5: 하단 카드 내 3개 영역으로 명확하게 구분해줘."
    )

    assert instructions == {
        3: "LangChain RAG Pipeline 아키텍처 다이어그램으로 변경해줘.",
        5: "하단 카드 내 3개 영역으로 명확하게 구분해줘.",
    }


def test_revision_request_expands_multi_slide_korean_targets() -> None:
    query = "슬라이드 4, 5 장 아키텍처 말고 일반 내용으로 변경해줘"

    assert _extract_slide_revision_instructions(query) == {
        4: "아키텍처 말고 일반 내용으로 변경해줘",
        5: "아키텍처 말고 일반 내용으로 변경해줘",
    }
    assert _normalize_revision_scope(None, query) == "slide_rewrite"


def test_negative_architecture_revision_suppresses_visual_asset_signal() -> None:
    query = "슬라이드 4, 5 장 아키텍처 말고 일반 내용으로 변경해줘"

    assert _negative_visual_asset_signal(query) is True
    assert _quick_visual_signal(
        query,
        [{"index": 4, "title": "Native Agent Architectures"}],
    ) is False


def test_blueprints_are_normalized_to_known_body_layouts() -> None:
    ruleset = get_ruleset()
    blueprints = _normalize_blueprints(
        [
            {"slide_type": "cover", "title": "Opening"},
            {
                "slide_type": "data",
                "title": "Evidence",
                "layout_plan": {
                    "body_layout_id": "not-a-layout",
                    "sub_layout_ids": ["grid_2x2", "not-a-layout"],
                },
            },
        ],
        ruleset,
    )

    assert blueprints[0]["layout_plan"]["master_role"] == "cover"
    assert blueprints[1]["layout_plan"]["body_layout_id"] == "dashboard_chart_sidebar"
    assert blueprints[1]["layout_plan"]["sub_layout_ids"] == ["grid_2x2"]


def test_architecture_slide_does_not_reserve_empty_visual_asset_slot_by_default() -> None:
    blueprints = _normalize_blueprints(
        [
            {
                "index": 2,
                "slide_type": "content",
                "title": "LangGraph service architecture",
                "suggested_elements": ["diagram", "connector", "card"],
                "layout_plan": {"body_layout_id": "split_60_40"},
            }
        ],
        get_ruleset(),
    )

    placements = blueprints[0]["layout_plan"]["element_placements"]
    assert not any(item.get("asset_role") == "visual_asset" for item in placements)
    assert any(item.get("element") == "card" for item in placements)
    assert len(placements) >= 5


def test_generic_flow_suggested_elements_are_demoted_without_explicit_visual_request() -> None:
    blueprints = _normalize_blueprints(
        [
            {
                "index": 2,
                "slide_type": "content",
                "title": "Native document generation flow",
                "suggested_elements": ["diagram", "flowchart", "image", "table"],
                "layout_plan": {"body_layout_id": "split_60_40"},
            }
        ],
        get_ruleset(),
        "제품 소개서를 만들어줘",
    )

    suggested = {item.lower() for item in blueprints[0]["suggested_elements"]}

    assert "diagram" not in suggested
    assert "flowchart" not in suggested
    assert "image" not in suggested
    assert {"connector", "rounded_rect", "table"} <= suggested


def test_explicit_redraw_request_reserves_visual_asset_slot() -> None:
    blueprints = _normalize_blueprints(
        [
            {
                "index": 2,
                "slide_type": "content",
                "title": "LangGraph service architecture",
                "suggested_elements": ["diagram", "connector", "card"],
                "layout_plan": {"body_layout_id": "split_60_40"},
            }
        ],
        get_ruleset(),
        "architecture diagram을 다시 그려줘",
    )

    placements = blueprints[0]["layout_plan"]["element_placements"]
    visual_slot = next(item for item in placements if item.get("asset_role") == "visual_asset")
    assert visual_slot["element"] == "image"
    assert visual_slot["w"] >= 560
    assert visual_slot["h"] >= 300
    assert len(placements) >= 5


def test_overlapping_explicit_visual_placements_repair_to_visual_slot() -> None:
    blueprints = _normalize_blueprints(
        [
            {
                "index": 2,
                "slide_type": "content",
                "title": "Agentic flow diagram",
                "layout_plan": {
                    "body_layout_id": "split_60_40",
                    "element_placements": [
                        {
                            "id": "diagram",
                            "element": "image",
                            "asset_role": "visual_asset",
                            "x": 100,
                            "y": 110,
                            "w": 560,
                            "h": 300,
                        },
                        {"id": "card_a", "element": "card", "x": 120, "y": 130, "w": 240, "h": 90},
                        {"id": "card_b", "element": "card", "x": 380, "y": 130, "w": 240, "h": 90},
                    ],
                },
            }
        ],
        get_ruleset(),
        "diagram redraw",
    )

    placements = blueprints[0]["layout_plan"]["element_placements"]

    assert any(item.get("asset_role") == "visual_asset" for item in placements)
    assert not _placements_overlap_for_test(placements)


def test_html_prompt_demotes_visual_asset_slot_when_no_rendered_asset_exists() -> None:
    placements = [
        {
            "id": "visual_asset_main",
            "element": "image",
            "role": "proof_object",
            "asset_role": "visual_asset",
            "x": 48,
            "y": 92,
            "w": 604,
            "h": 318,
            "fit": "contain",
        }
    ]

    prompt_placements = _element_placements_for_html_prompt(
        placements,
        has_visual_assets=False,
    )

    assert prompt_placements[0]["element"] == "card"
    assert prompt_placements[0]["role"] == "proof_object"
    assert "asset_role" not in prompt_placements[0]
    assert "content_requirement" in prompt_placements[0]


def test_requested_slide_count_is_extended_when_planner_under_produces() -> None:
    ruleset = get_ruleset()
    blueprints = _normalize_blueprints(
        [
            {"index": 1, "slide_type": "cover", "title": "Opening"},
            {"index": 2, "slide_type": "content", "title": "Problem"},
            {"index": 3, "slide_type": "content", "title": "Solution"},
        ],
        ruleset,
    )

    extended = _extend_blueprints_to_requested_count(
        blueprints,
        6,
        "6 slides product brief",
        ruleset,
    )

    assert len(extended) == 6
    assert extended[-1]["slide_type"] == "summary"
    assert extended[-1]["layout_plan"]["element_placements"]


def test_requested_slide_count_extension_uses_user_source_material() -> None:
    ruleset = get_ruleset()
    blueprints = _normalize_blueprints(
        [
            {"index": 1, "slide_type": "cover", "title": "DocuMind"},
            {"index": 2, "slide_type": "content", "title": "제품 한 줄"},
            {"index": 3, "slide_type": "content", "title": "지원 포맷"},
        ],
        ruleset,
    )
    query = """
    # DocuMind 상품 소개서 근거
    ## 1. 제품 한 줄
    - LangGraph 기반 Agentic AI 문서 자동화 플랫폼
    ## 2. 지원 포맷
    - PPTX, DOCX, PDF, Markdown, XLSX, HWPX 네이티브 생성
    ## 3. 소비 채널
    - FastAPI, Web UI, PyPI 패키지 3면 소비
    ## 4. 오케스트레이션
    - init, plan, visual asset, generate HTML, render, QA 단계
    ## 5. 품질 보증
    - deterministic OOXML mapper와 VLM QA 루프를 결합
    """

    extended = _extend_blueprints_to_requested_count(
        blueprints,
        7,
        query,
        ruleset,
    )

    assert len(extended) == 7
    assert "소비 채널" in extended[3]["title"]
    assert "오케스트레이션" in extended[4]["title"]
    assert "품질 보증" in extended[5]["title"]
    assert "deterministic OOXML mapper" in str(extended[5]["content_blocks"])
    assert "핵심 설계 포인트" not in str(extended)
    assert "보강 슬라이드" not in str(extended)


def test_compact_parent_slide_plan_preserves_late_slide_context() -> None:
    slide_plan = [
        {
            "index": index,
            "slide_type": "content",
            "title": f"Slide {index} source-backed title",
            "key_message": f"Key message {index}",
            "content_blocks": [
                {
                    "type": "points",
                    "items": [
                        {"title": f"Finding {index}", "body": "x" * 600},
                    ],
                }
            ],
        }
        for index in range(1, 14)
    ]

    compact = _compact_parent_slide_plan(slide_plan, max_chars=6000)

    assert "Slide 1 source-backed title" in compact
    assert "Slide 13 source-backed title" in compact


def test_slide_outline_scope_ignores_p99_latency_metric() -> None:
    query = """
    Slide 1 — 운영 관점
    핵심 메시지: latency, error rate, throughput을 본다.
    Slide 2 — 5대 관측 축
    Metrics
    QPS, P99 latency, 토큰, 비용
    Slide 3 — 정리
    """

    explicit_count = _extract_requested_slide_count(query)
    instructions = _extract_slide_revision_instructions(query)

    assert explicit_count is None
    assert sorted(instructions) == [1, 2, 3]
    assert _infer_target_slide_count(
        query,
        explicit_slide_count=explicit_count,
        base_version={},
        explicit_slide_instructions=instructions,
    ) == 3


async def test_missing_slide_count_repair_requests_only_missing_slides() -> None:
    ruleset = get_ruleset()
    existing = _normalize_blueprints(
        [
            {"index": 1, "slide_type": "cover", "title": "Opening"},
            {"index": 2, "slide_type": "content", "title": "Problem"},
            {"index": 3, "slide_type": "content", "title": "Solution"},
        ],
        ruleset,
    )

    class FakeResponse:
        content = """
        {
          "slides": [
            {
              "index": 4,
              "slide_type": "content",
              "title": "Execution Plan",
              "key_message": "Pilot the workflow before scaling.",
              "purpose": "Show the rollout path.",
              "content_blocks": [
                {"type": "steps", "items": [{"title": "Pilot", "body": "Start with one team."}]}
              ],
              "layout_plan": {"body_layout_id": "process_4col"},
              "suggested_elements": ["rounded_rect", "connector"]
            },
            {
              "index": 5,
              "slide_type": "summary",
              "title": "Decision Summary",
              "key_message": "The requested deck now has the complete arc.",
              "purpose": "Close with the decision.",
              "content_blocks": [
                {"type": "summary", "items": [{"title": "Approve", "body": "Move to pilot."}]}
              ],
              "layout_plan": {"body_layout_id": "numbered_list_card"},
              "suggested_elements": ["rounded_rect", "icon"]
            }
          ]
        }
        """

    class FakeLLM:
        def __init__(self) -> None:
            self.prompt = ""

        async def ainvoke(self, messages):
            self.prompt = messages[-1].content
            return FakeResponse()

    llm = FakeLLM()
    result = await _repair_missing_slide_count_result(
        llm=llm,
        base_messages=[],
        parsed_result={"title": "Deck", "slides": existing},
        slide_blueprints=existing,
        requested_count=5,
        user_query="5 slides product brief",
        ruleset=ruleset,
        research_data=None,
    )

    assert result is not None
    assert len(result["slides"]) == 5
    assert [slide["index"] for slide in result["slides"]] == [1, 2, 3, 4, 5]
    assert "Do NOT rewrite the whole deck" in llm.prompt
    assert "containing exactly the 2 missing slide blueprints" in llm.prompt


def test_visual_variant_request_infers_two_slide_scope() -> None:
    query = (
        "첨부된 이미지의 아키텍처에 트랜짓 게이트웨이를 추가해서 다시 그려줘. "
        "슬라이드 한장은 다이어그램스로 한장은 이미지 모델로 그려줘."
    )

    assert _infer_target_slide_count(
        query,
        explicit_slide_count=None,
        base_version={},
        explicit_slide_instructions={},
    ) == 2


def test_existing_revision_defaults_to_parent_slide_count() -> None:
    assert _infer_target_slide_count(
        "아키텍처 다이어그램을 다시 그려줘",
        explicit_slide_count=None,
        base_version={"slide_plan": [{"index": 1}, {"index": 2}, {"index": 3}]},
        explicit_slide_instructions={},
    ) == 3


def test_repeated_missing_layouts_are_diversified_by_position() -> None:
    blueprints = _normalize_blueprints(
        [
            {"index": index, "slide_type": "content", "title": f"Architecture {index}"}
            for index in range(1, 6)
        ],
        get_ruleset(),
    )

    layout_ids = [bp["layout_plan"]["body_layout_id"] for bp in blueprints]

    assert len(set(layout_ids)) >= 4
    assert len({layout_id.split("_", 1)[0] for layout_id in layout_ids}) >= 3


def test_valid_llm_selected_layouts_are_not_rewritten_for_diversity() -> None:
    blueprints = _normalize_blueprints(
        [
            {
                "index": index,
                "slide_type": "content",
                "title": f"Comparison {index}",
                "layout_plan": {"body_layout_id": "compare_vs"},
            }
            for index in range(1, 5)
        ],
        get_ruleset(),
    )

    assert [bp["layout_plan"]["body_layout_id"] for bp in blueprints] == [
        "compare_vs",
        "compare_vs",
        "compare_vs",
        "compare_vs",
    ]


def test_llm_json_parser_uses_first_complete_object_when_extra_data_follows() -> None:
    parsed = parse_llm_json(
        '{"score": 0.91, "passed": true, "issues": []}\n'
        '{"score": 0.1, "passed": false}'
    )

    assert parsed == {"score": 0.91, "passed": True, "issues": []}


async def test_unified_planner_prefers_structured_output_when_available(monkeypatch) -> None:
    from src.core.config import settings

    monkeypatch.setattr(settings, "llm_provider", "openai")

    class FakeStructuredRunnable:
        async def ainvoke(self, messages):
            return {
                "title": "Structured plan",
                "slides": [{"index": 1, "slide_type": "cover", "title": "Opening"}],
                "design_tokens": {"primary": "#111827"},
            }

    class FakeLLM:
        def __init__(self) -> None:
            self.schema = None
            self.kwargs = None

        def with_structured_output(self, schema, **kwargs):
            self.schema = schema
            self.kwargs = kwargs
            return FakeStructuredRunnable()

    llm = FakeLLM()
    result = await _invoke_planner_with_structured_output(
        llm=llm,
        messages=[],
        requested_count=1,
    )

    assert result["title"] == "Structured plan"
    assert result["slides"][0]["title"] == "Opening"
    assert llm.kwargs == {"method": "json_schema"}


def test_unified_planner_skips_structured_output_for_bedrock(monkeypatch) -> None:
    from src.core.config import settings

    class FakeLLM:
        def with_structured_output(self, schema, **kwargs):
            raise AssertionError("Bedrock should skip structured output")

    monkeypatch.setattr(settings, "llm_provider", "bedrock")

    assert _structured_planner_llm(FakeLLM()) is None


async def test_unified_planner_rejects_empty_structured_output() -> None:
    class FakeStructuredRunnable:
        async def ainvoke(self, messages):
            return {"title": "Incomplete plan"}

    class FakeLLM:
        def with_structured_output(self, schema, **kwargs):
            return FakeStructuredRunnable()

    result = await _invoke_planner_with_structured_output(
        llm=FakeLLM(),
        messages=[],
        requested_count=None,
    )

    assert result is None


async def test_nova_canvas_image_generation_uses_bedrock_text_image_payload(monkeypatch) -> None:
    import base64
    import json

    captured = {}

    class FakeBody:
        def read(self):
            return json.dumps({"images": [base64.b64encode(b"png").decode("ascii")]}).encode()

    class FakeClient:
        def invoke_model(self, **kwargs):
            captured.update(kwargs)
            return {"body": FakeBody()}

    monkeypatch.setattr("src.utils.image_gen._bedrock_runtime_client", lambda: FakeClient())

    image = await _generate_with_nova_canvas(
        "amazon.nova-canvas-v1:0",
        "Clean product illustration",
        "watermark",
        1024,
        1024,
    )

    payload = json.loads(captured["body"])
    assert captured["modelId"] == "amazon.nova-canvas-v1:0"
    assert payload["taskType"] == "TEXT_IMAGE"
    assert payload["textToImageParams"]["text"] == "Clean product illustration"
    assert payload["textToImageParams"]["negativeText"] == "watermark"
    assert payload["imageGenerationConfig"]["width"] == 1024
    assert image == b"png"


def test_overlapping_geometric_placements_are_replaced_with_safe_layout() -> None:
    blueprints = _normalize_blueprints(
        [
            {
                "index": 2,
                "slide_type": "content",
                "title": "Dense layout",
                "layout_plan": {
                    "body_layout_id": "split_60_40",
                    "element_placements": [
                        {"id": "a", "element": "card", "x": 40, "y": 100, "w": 400, "h": 240},
                        {"id": "b", "element": "card", "x": 80, "y": 120, "w": 400, "h": 240},
                        {"id": "c", "element": "callout", "x": 40, "y": 430, "w": 880, "h": 70},
                    ],
                },
            }
        ],
        get_ruleset(),
    )

    placements = blueprints[0]["layout_plan"]["element_placements"]
    assert placements[0]["id"] == "main_proof"
    assert placements[1]["x"] >= placements[0]["x"] + placements[0]["w"]


def test_fixed_template_uses_selected_master_zone_positions() -> None:
    master_layout = get_ruleset().resolve_master_layout(
        {"header_zone_id": "header_side_rule", "footer_zone_id": "footer_compact"}
    )
    html = (
        '<div data-slide="2" style="position:absolute;left:0;top:0;width:960px;height:540px">'
        '<div data-pptx-type="textbox" style="position:absolute;left:40px;top:100px;'
        'width:200px;height:40px">Body</div></div>'
    )
    output = _inject_fixed_template(
        html,
        2,
        {"title": "A decision claim", "section_label": "EVIDENCE"},
        {"master_layout": master_layout},
    )

    assert "left:52px;top:16px;width:250px;height:14px" in output
    assert "left:40px;top:16px;width:4px;height:54px" in output
    assert "left:40px;top:526px;width:720px;height:11px" in output
    assert ">Body</div>" in output


def test_header_title_size_is_consistent_for_long_and_short_titles() -> None:
    short = _inject_fixed_template(
        '<div data-slide="1"><div data-pptx-type="textbox" '
        'style="position:absolute;left:40px;top:100px;width:200px;height:40px">Body</div></div>',
        1,
        {"title": "Short title", "section_label": "A"},
        {},
    )
    long = _inject_fixed_template(
        '<div data-slide="2"><div data-pptx-type="textbox" '
        'style="position:absolute;left:40px;top:100px;width:200px;height:40px">Body</div></div>',
        2,
        {
            "title": "AI 관측 서비스 구축을 위한 전략적 실행 로드맵과 핵심 과제",
            "section_label": "A",
        },
        {},
    )

    assert "font-size:22px" in short
    assert "font-size:22px" in long
    assert "line-height:1.08;padding:0;vertical-align:middle" in long


def test_template_slide_size_scales_html_coordinates(tmp_path: Path) -> None:
    template = Presentation()
    template.slide_width = Emu(12_192_000)
    template.slide_height = Emu(6_858_000)
    template.slides.add_slide(template.slide_layouts[6])
    buffer = BytesIO()
    template.save(buffer)

    html = (
        '<div data-slide="1">'
        '<div data-pptx-type="textbox" style="position:absolute;left:480px;top:270px;'
        'width:96px;height:54px;font-size:20px;color:#111827">Scaled</div></div>'
    )

    output = CSStoOOXMLEngine().build_presentation(
        [{"index": 1, "html": html, "metadata": {"slide_type": "content"}}],
        tmp_path,
        template_bytes=buffer.getvalue(),
    )
    prs = Presentation(output)
    shape = prs.slides[0].shapes[0]

    assert prs.slide_width == 12_192_000
    assert prs.slide_height == 6_858_000
    assert abs(shape.left - 6_096_000) < 5
    assert abs(shape.top - 3_429_000) < 5
    assert abs(shape.width - 1_219_200) < 5


def test_icon_is_layered_above_card_without_moving_card(tmp_path: Path, monkeypatch) -> None:
    icon_path = tmp_path / "icon.png"
    Image.new("RGBA", (32, 32), (12, 34, 56, 255)).save(icon_path)
    monkeypatch.setattr("src.utils.iconify.get_icon_asset_path", lambda *args, **kwargs: icon_path)

    element = ParsedElement(
        pptx_type="textbox",
        pptx_shape=None,
        position={"left": 40, "top": 100, "width": 220, "height": 120},
        styles={
            "background-color": "#F1F5F9",
            "color": "#112233",
            "padding": "12px",
        },
        text_content="Data pipeline",
        children=[],
        attributes={"data-pptx-icon": "database"},
    )
    original_position = dict(element.position)
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])

    CSStoOOXMLEngine()._add_element(slide, element)

    assert element.position == original_position
    assert element.styles["padding-top"] == "50px"
    assert slide.shapes[0].shape_type != MSO_SHAPE_TYPE.PICTURE
    assert slide.shapes[1].shape_type == MSO_SHAPE_TYPE.PICTURE


def test_explicit_icon_slot_renders_on_compact_heading(tmp_path: Path, monkeypatch) -> None:
    icon_path = tmp_path / "icon.png"
    Image.new("RGBA", (32, 32), (12, 34, 56, 255)).save(icon_path)
    monkeypatch.setattr("src.utils.iconify.get_icon_asset_path", lambda *args, **kwargs: icon_path)

    element = ParsedElement(
        pptx_type="textbox",
        pptx_shape=None,
        position={"left": 40, "top": 100, "width": 220, "height": 24},
        styles={"color": "#112233", "font-size": "14px"},
        text_content="Compact heading",
        children=[],
        attributes={
            "data-pptx-icon": "database",
            "data-pptx-icon-layout": "inline-left",
            "data-pptx-icon-size": "20",
        },
    )
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])

    CSStoOOXMLEngine()._add_element(slide, element)

    assert element.styles["padding-left"] == "34px"
    assert slide.shapes[0].shape_type != MSO_SHAPE_TYPE.PICTURE
    assert slide.shapes[1].shape_type == MSO_SHAPE_TYPE.PICTURE


def test_html_preview_embeds_icons_for_explicit_compact_slots(monkeypatch, tmp_path: Path) -> None:
    icon_path = tmp_path / "icon.svg"
    icon_path.write_text(
        '<svg xmlns="http://www.w3.org/2000/svg" width="32" height="32">'
        '<path fill="#112233" d="M4 4h24v24H4z"/></svg>',
        encoding="utf-8",
    )
    monkeypatch.setattr("src.utils.iconify.get_icon_asset_path", lambda *args, **kwargs: icon_path)

    html = (
        '<div data-slide="1">'
        '<div data-pptx-type="textbox" data-pptx-icon="database" '
        'data-pptx-icon-layout="inline-left" data-pptx-icon-size="20" '
        'style="position:absolute;left:40px;top:100px;width:220px;height:24px;'
        'color:#112233;font-size:14px">Compact heading</div></div>'
    )

    output = _embed_cached_icons(html)

    assert "data:image/svg+xml;base64" in output
    assert "width:20px;height:20px" in output
    assert "padding-left:34px" in output


def test_independent_icon_element_maps_to_html_and_pptx(tmp_path: Path, monkeypatch) -> None:
    png_path = tmp_path / "icon.png"
    svg_path = tmp_path / "icon.svg"
    Image.new("RGBA", (32, 32), (12, 34, 56, 255)).save(png_path)
    svg_path.write_text(
        '<svg xmlns="http://www.w3.org/2000/svg" width="32" height="32">'
        '<path fill="#112233" d="M4 4h24v24H4z"/></svg>',
        encoding="utf-8",
    )

    def fake_asset(*args, **kwargs):
        return svg_path if kwargs.get("target") == "html" else png_path

    monkeypatch.setattr("src.utils.iconify.get_icon_asset_path", fake_asset)
    html = (
        '<div data-slide="1">'
        '<div data-pptx-type="icon" data-pptx-icon="brain" '
        'data-pptx-icon-placement="diagram_node_top" '
        'style="position:absolute;left:40px;top:100px;width:32px;height:32px;'
        'color:#112233"></div>'
        '<div data-pptx-type="textbox" style="position:absolute;left:84px;top:96px;'
        'width:220px;height:44px;color:#112233">AI Engine</div></div>'
    )

    preview = _embed_cached_icons(html)
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    for element in parse_slide_html(html):
        CSStoOOXMLEngine()._add_element(slide, element)

    assert "data:image/svg+xml;base64" in preview
    assert "left:40px;top:100px;width:32px;height:32px" in preview
    assert slide.shapes[0].shape_type == MSO_SHAPE_TYPE.PICTURE
    assert slide.shapes[1].shape_type != MSO_SHAPE_TYPE.PICTURE


def test_independent_icon_is_centered_in_declared_box(tmp_path: Path, monkeypatch) -> None:
    icon_path = tmp_path / "wide_icon.png"
    Image.new("RGBA", (64, 32), (12, 34, 56, 255)).save(icon_path)
    monkeypatch.setattr("src.utils.iconify.get_icon_asset_path", lambda *args, **kwargs: icon_path)

    html = (
        '<div data-slide="1">'
        '<div data-pptx-type="icon" data-pptx-icon="database" '
        'style="position:absolute;left:40px;top:100px;width:40px;height:40px;'
        'color:#112233"></div></div>'
    )
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    for element in parse_slide_html(html):
        CSStoOOXMLEngine()._add_element(slide, element)

    picture = slide.shapes[0]
    assert picture.shape_type == MSO_SHAPE_TYPE.PICTURE
    assert abs(picture.left - Emu(40 * 9525)) < 5
    assert abs(picture.top - Emu(110 * 9525)) < 5
    assert abs(picture.width - Emu(40 * 9525)) < 5
    assert abs(picture.height - Emu(20 * 9525)) < 5


def test_generated_html_icon_node_is_materialized_with_data_image(tmp_path: Path) -> None:
    from bs4 import BeautifulSoup

    icon_path = tmp_path / "icon.svg"
    icon_path.write_text(
        '<svg xmlns="http://www.w3.org/2000/svg" width="32" height="32">'
        '<path fill="#112233" d="M4 4h24v24H4z"/></svg>',
        encoding="utf-8",
    )
    soup = BeautifulSoup(
        '<div data-pptx-type="icon" data-pptx-icon="brain" '
        'style="position:absolute;left:40px;top:100px;width:32px;height:32px;'
        'background-color:#A855F7;color:#112233"></div>',
        "html.parser",
    )

    _materialize_html_icon_node(soup.div, icon_path)

    style = soup.div["style"]
    assert "background-image:url(data:image/svg+xml;base64," in style
    assert "background-color:transparent" in style
    assert "background-color:#A855F7" not in style


def test_legacy_textbox_icon_is_normalized_to_separate_icon_area() -> None:
    html = (
        '<div data-slide="1">'
        '<div data-pptx-type="textbox" data-pptx-icon="brain" '
        'style="position:absolute;left:40px;top:100px;width:220px;height:32px;'
        'color:#112233;font-size:14px">AI Engine</div></div>'
    )

    output = _normalize_legacy_icon_nodes(html)

    assert 'data-pptx-type="icon"' in output
    assert 'data-pptx-icon-placement="card_lead_left"' in output
    assert output.count('data-pptx-icon="brain"') == 1
    assert "left:74px" in output


def test_parser_recovers_flattened_bullet_line_breaks() -> None:
    html = (
        '<div data-slide="1">'
        '<div data-pptx-type="textbox" style="position:absolute;left:0;top:0;'
        'width:240px;height:120px">• 클라우드 MSP 전문 기업 • 토스뱅크 주주 참여 '
        '• 2024년 매출 성장</div></div>'
    )

    [element] = parse_slide_html(html)

    assert element.text_content.splitlines() == [
        "• 클라우드 MSP 전문 기업",
        "• 토스뱅크 주주 참여",
        "• 2024년 매출 성장",
    ]


def test_html_lists_are_preserved_as_bullet_lines() -> None:
    html = (
        '<div data-slide="1">'
        '<div data-pptx-type="textbox" style="position:absolute;left:40px;top:100px;'
        'width:240px;height:120px;font-size:14px">'
        "<ul><li>Data collection</li><li>Model training</li><li>Deployment</li></ul>"
        "</div></div>"
    )

    normalized = _normalize_slide_html(html)
    [element] = parse_slide_html(normalized)

    assert 'data-pptx-list="bullet"' in normalized
    assert element.text_content.splitlines() == [
        "• Data collection",
        "• Model training",
        "• Deployment",
    ]


def test_inline_bullet_text_is_materialized_as_html_list() -> None:
    html = (
        '<div data-slide="1">'
        '<div data-pptx-type="textbox" style="position:absolute;left:40px;top:100px;'
        'width:260px;height:72px;font-size:14px">'
        "• Data collection • Model training • Deployment"
        "</div></div>"
    )

    normalized = _normalize_slide_html(html)
    [element] = parse_slide_html(normalized)

    assert "<ul>" in normalized
    assert normalized.count("<li>") == 3
    assert 'data-pptx-list="bullet"' in normalized
    assert element.text_content.splitlines() == [
        "• Data collection",
        "• Model training",
        "• Deployment",
    ]


def test_list_normalization_removes_duplicate_manual_bullets() -> None:
    html = (
        '<div data-slide="1">'
        '<div data-pptx-type="textbox" style="position:absolute;left:40px;top:100px;'
        'width:260px;height:72px;font-size:14px">'
        "<ul><li>\u2022 Data collection</li><li>\u2022 \u2022 Model training</li></ul>"
        "</div></div>"
    )

    normalized = _normalize_slide_html(html)
    [element] = parse_slide_html(normalized)
    pptx_text = CSStoOOXMLEngine()._normalize_list_text(
        "\u2022 \u2022 Data collection\n\u2022 Model training",
        ordered=False,
    )

    assert normalized.count("<li>") == 2
    assert "\u2022 \u2022" not in normalized
    assert element.text_content.splitlines() == [
        "\u2022 Data collection",
        "\u2022 Model training",
    ]
    assert pptx_text.splitlines() == [
        "\u2022 Data collection",
        "\u2022 Model training",
    ]


def test_inline_flow_arrow_text_is_not_split_as_list() -> None:
    html = (
        '<div data-slide="1">'
        '<div data-pptx-type="textbox" style="position:absolute;left:40px;top:100px;'
        'width:320px;height:40px;font-size:14px">'
        "Input → Engine → Output"
        "</div></div>"
    )

    [element] = parse_slide_html(html)

    assert element.text_content.splitlines() == ["Input → Engine → Output"]


def test_list_textbox_and_backing_card_expand_to_fit_lines() -> None:
    html = (
        '<div data-slide="1">'
        '<div data-pptx-type="shape" data-pptx-shape="rounded_rect" '
        'style="position:absolute;left:40px;top:100px;width:260px;height:56px;'
        'background-color:#FFFFFF;border:1px solid #CBD5E1"></div>'
        '<div data-pptx-type="textbox" style="position:absolute;left:52px;top:112px;'
        'width:236px;height:28px;font-size:14px;line-height:1.5;color:#111827">'
        "<ul><li>Data collection</li><li>Model training</li><li>Deployment</li></ul>"
        "</div></div>"
    )

    normalized = _normalize_slide_html(html)
    elements = parse_slide_html(normalized)
    card, textbox = elements[0], elements[1]

    assert textbox.position["height"] > 28
    assert card.position["height"] >= (
        textbox.position["top"] + textbox.position["height"] - card.position["top"] + 8
    )
    assert card.position["height"] <= 526 - card.position["top"]


def test_content_boxes_are_clamped_above_footer_safe_area() -> None:
    html = (
        '<div data-slide="1" style="position:absolute;left:0;top:0;width:960px;height:540px">'
        '<div data-pptx-type="shape" data-pptx-shape="rounded_rect" '
        'style="position:absolute;left:40px;top:480px;width:880px;height:80px;'
        'background-color:#ECFDF5"></div>'
        '<div data-pptx-type="textbox" style="position:absolute;left:60px;top:492px;'
        'width:820px;height:48px;font-size:13px;color:#064E3B">'
        "Footer-safe callout text"
        "</div></div>"
    )

    elements = parse_slide_html(_normalize_slide_html(html))

    for element in elements:
        assert element.position["top"] + element.position["height"] <= 510


def test_overlapping_card_groups_move_with_their_icons_and_text() -> None:
    html = (
        '<div data-slide="1" style="position:absolute;left:0;top:0;width:960px;height:540px">'
        '<div data-pptx-type="shape" data-pptx-shape="rounded_rect" '
        'style="position:absolute;left:40px;top:100px;width:220px;height:120px;'
        'background-color:#ECFDF5"></div>'
        '<div data-pptx-type="icon" data-pptx-icon="database" '
        'style="position:absolute;left:52px;top:112px;width:24px;height:24px;'
        'color:#064E3B"></div>'
        '<div data-pptx-type="textbox" style="position:absolute;left:88px;top:110px;'
        'width:160px;height:32px;font-size:14px;font-weight:700;color:#064E3B">Card A</div>'
        '<div data-pptx-type="shape" data-pptx-shape="rounded_rect" '
        'style="position:absolute;left:60px;top:120px;width:220px;height:120px;'
        'background-color:#DBEAFE"></div>'
        '<div data-pptx-type="icon" data-pptx-icon="brain" '
        'style="position:absolute;left:72px;top:132px;width:24px;height:24px;'
        'color:#1E3A8A"></div>'
        '<div data-pptx-type="textbox" style="position:absolute;left:108px;top:130px;'
        'width:160px;height:32px;font-size:14px;font-weight:700;color:#1E3A8A">Card B</div>'
        "</div>"
    )

    elements = parse_slide_html(_normalize_slide_html(html))
    card_a = elements[0]
    card_b = elements[3]
    card_b_text = next(element for element in elements if element.text_content == "Card B")

    horizontal_gap = card_b.position["left"] - (card_a.position["left"] + card_a.position["width"])
    vertical_gap = card_b.position["top"] - (card_a.position["top"] + card_a.position["height"])
    assert horizontal_gap >= 0 or vertical_gap >= 0
    assert card_b_text.position["top"] > 130 or card_b_text.position["left"] > 108


def test_large_overlapping_card_is_not_treated_as_inner_decoration() -> None:
    html = (
        '<div data-slide="1" style="position:absolute;left:0;top:0;width:960px;height:540px">'
        '<div data-pptx-type="shape" data-pptx-shape="rounded_rect" '
        'style="position:absolute;left:100px;top:120px;width:420px;height:300px;'
        'background-color:#111827"></div>'
        '<div data-pptx-type="textbox" style="position:absolute;left:124px;top:142px;'
        'width:340px;height:36px;font-size:18px;color:#FFFFFF">Left card</div>'
        '<div data-pptx-type="shape" data-pptx-shape="rounded_rect" '
        'style="position:absolute;left:360px;top:120px;width:268px;height:300px;'
        'background-color:#0F766E"></div>'
        '<div data-pptx-type="textbox" style="position:absolute;left:384px;top:142px;'
        'width:210px;height:36px;font-size:18px;color:#FFFFFF">Right card</div>'
        "</div>"
    )

    elements = parse_slide_html(_normalize_slide_html(html))
    left_card = elements[0]
    right_card = elements[2]

    assert right_card.position["left"] >= left_card.position["left"] + left_card.position["width"]


def test_empty_large_card_container_is_removed_before_rendering() -> None:
    html = (
        '<div data-slide="1" style="position:absolute;left:0;top:0;width:960px;height:540px">'
        '<div data-pptx-type="shape" data-pptx-shape="rounded_rect" '
        'style="position:absolute;left:80px;top:120px;width:320px;height:180px;'
        'background-color:#CCFBF1"></div>'
        '<div data-pptx-type="shape" data-pptx-shape="rounded_rect" '
        'style="position:absolute;left:460px;top:120px;width:320px;height:180px;'
        'background-color:#DBEAFE"></div>'
        '<div data-pptx-type="textbox" style="position:absolute;left:488px;top:148px;'
        'width:250px;height:40px;font-size:18px;color:#1E3A8A">Filled card</div>'
        "</div>"
    )

    elements = parse_slide_html(_normalize_slide_html(html))

    assert len([element for element in elements if element.pptx_shape == "rounded_rect"]) == 1
    assert any(element.text_content == "Filled card" for element in elements)


def test_layout_inset_shrinks_slots_without_replanning_card_gaps() -> None:
    placements = [
        {"id": "card_a", "element": "card", "x": 40, "y": 120, "w": 180, "h": 100},
        {"id": "card_b", "element": "card", "x": 226, "y": 120, "w": 180, "h": 100},
    ]

    inset = _apply_layout_inset(placements)

    assert inset[0] == {"id": "card_a", "element": "card", "x": 44, "y": 124, "w": 172, "h": 92}
    assert inset[1] == {"id": "card_b", "element": "card", "x": 230, "y": 124, "w": 172, "h": 92}
    assert placements[1]["x"] - (placements[0]["x"] + placements[0]["w"]) == 6


def test_parser_inherits_root_and_inline_font_details() -> None:
    html = (
        '<div data-slide="1">'
        '<div data-pptx-type="textbox" style="position:absolute;left:40px;top:100px;'
        "width:300px;height:80px;font-family:'Inter';font-size:18px;"
        'font-weight:600;color:#112233">'
        'Root <span style="font-size:12px;color:#CC0000;font-style:italic;'
        'text-decoration:underline">inline</span></div></div>'
    )

    [element] = parse_slide_html(html)

    assert element.text_runs[0]["font_family"] == "'Inter'"
    assert element.text_runs[0]["size"] == "18px"
    assert element.text_runs[0]["bold"] is True
    inline = element.text_runs[1]
    assert inline["text"] == "inline"
    assert inline["size"] == "12px"
    assert inline["color"] == "#CC0000"
    assert inline["italic"] is True
    assert inline["underline"] is True


def test_normalized_html_forces_text_contrast_on_backing_shape() -> None:
    html = (
        '<div data-slide="1" style="position:absolute;left:0;top:0;width:960px;height:540px">'
        '<div data-pptx-type="shape" data-pptx-shape="rounded_rect" '
        'style="position:absolute;left:40px;top:100px;width:320px;height:120px;'
        'background-color:#111827"></div>'
        '<div data-pptx-type="textbox" style="position:absolute;left:60px;top:120px;'
        'width:280px;height:40px;font-size:14px;font-weight:400;color:#111827">'
        'Readable text</div></div>'
    )

    normalized = _normalize_slide_html(html)

    assert "color:#FFFFFF" in normalized


def test_normalized_html_uses_same_family_dark_text_on_light_fill() -> None:
    html = (
        '<div data-slide="1" style="position:absolute;left:0;top:0;width:960px;height:540px">'
        '<div data-pptx-type="textbox" style="position:absolute;left:40px;top:100px;'
        'width:280px;height:80px;background-color:#ECFDF5;font-size:14px;'
        'font-weight:400;color:#FFFFFF">Light green card</div></div>'
    )

    normalized = _normalize_slide_html(html)

    assert "color:#FFFFFF" not in normalized
    assert "color:#073620" in normalized


def test_normalized_html_forces_inline_span_contrast_on_dark_fill() -> None:
    html = (
        '<div data-slide="1" style="position:absolute;left:0;top:0;width:960px;height:540px">'
        '<div data-pptx-type="shape" data-pptx-shape="rounded_rect" '
        'style="position:absolute;left:40px;top:100px;width:320px;height:120px;'
        'background-color:#111827"></div>'
        '<div data-pptx-type="textbox" style="position:absolute;left:60px;top:120px;'
        'width:280px;height:40px;font-size:14px;color:#FFFFFF">'
        '<span style="color:#111827">Inline dark text</span></div></div>'
    )

    normalized = _normalize_slide_html(html)

    assert 'span style="color:#FFFFFF"' in normalized


def test_normalized_html_uses_same_family_dark_text_on_light_backing_shape() -> None:
    html = (
        '<div data-slide="1" style="position:absolute;left:0;top:0;width:960px;height:540px">'
        '<div data-pptx-type="shape" data-pptx-shape="rounded_rect" '
        'style="position:absolute;left:40px;top:100px;width:320px;height:120px;'
        'background-color:#FCE7F3"></div>'
        '<div data-pptx-type="textbox" style="position:absolute;left:60px;top:120px;'
        'width:280px;height:40px;font-size:14px;color:#FFFFFF">'
        'Readable text</div></div>'
    )

    normalized = _normalize_slide_html(html)

    assert "color:#FFFFFF" not in normalized
    assert "color:#360722" in normalized


def test_pptx_text_color_is_corrected_against_own_dark_fill() -> None:
    element = ParsedElement(
        pptx_type="textbox",
        pptx_shape=None,
        position={"left": 40, "top": 100, "width": 220, "height": 80},
        styles={
            "background-color": "#111827",
            "font-size": "14px",
            "color": "#111827",
        },
        text_content="Readable",
        children=[],
        attributes={},
    )
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])

    CSStoOOXMLEngine()._add_element(slide, element)
    run = slide.shapes[0].text_frame.paragraphs[0].runs[0]

    assert str(run.font.color.rgb) == "FFFFFF"


def test_pptx_text_color_is_corrected_to_same_family_dark_on_light_fill() -> None:
    element = ParsedElement(
        pptx_type="textbox",
        pptx_shape=None,
        position={"left": 40, "top": 100, "width": 220, "height": 80},
        styles={
            "background-color": "#ECFDF5",
            "font-size": "14px",
            "color": "#FFFFFF",
        },
        text_content="Readable",
        children=[],
        attributes={},
    )
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])

    CSStoOOXMLEngine()._add_element(slide, element)
    run = slide.shapes[0].text_frame.paragraphs[0].runs[0]

    assert str(run.font.color.rgb) == "073620"


def test_design_evaluator_flags_text_contrast_against_backing_fill() -> None:
    html = (
        '<div data-slide="1" style="position:absolute;left:0;top:0;width:960px;height:540px">'
        '<div data-pptx-type="shape" data-pptx-shape="rounded_rect" '
        'style="position:absolute;left:40px;top:100px;width:320px;height:120px;'
        'background-color:#111827"></div>'
        '<div data-pptx-type="textbox" style="position:absolute;left:60px;top:120px;'
        'width:280px;height:40px;font-size:14px;color:#111827">Unreadable</div></div>'
    )

    result = DesignQualityEvaluator(get_ruleset()).evaluate([{"index": 1, "html": html}], {})

    assert any(
        "Text contrast failed" in issue["message"]
        for issue in result.per_slide[0]["issues"]
    )


def test_design_evaluator_flags_empty_large_container() -> None:
    html = (
        '<div data-slide="1" style="position:absolute;left:0;top:0;width:960px;height:540px">'
        '<div data-pptx-type="shape" data-pptx-shape="rounded_rect" '
        'style="position:absolute;left:60px;top:120px;width:520px;height:280px;'
        'background-color:#FFFFFF;border:1px solid #CBD5E1;box-shadow:0 2px 6px #CBD5E1">'
        '</div>'
        '<div data-pptx-type="textbox" style="position:absolute;left:640px;top:130px;'
        'width:220px;height:40px;font-size:18px;font-weight:700;color:#111827">Side content</div>'
        '</div>'
    )

    result = DesignQualityEvaluator(get_ruleset()).evaluate([{"index": 1, "html": html}], {})

    assert any(
        "Empty visible container detected" in issue["message"]
        for issue in result.per_slide[0]["issues"]
    )


def test_design_evaluator_flags_plain_typography() -> None:
    textboxes = "".join(
        f'<div data-pptx-type="textbox" style="position:absolute;left:60px;top:{80 + i * 52}px;'
        'width:520px;height:34px;font-family:Pretendard;font-size:16px;'
        'font-weight:400;line-height:1.35;color:#111827">Item {i}</div>'
        for i in range(4)
    )
    html = (
        '<div data-slide="1" style="position:absolute;left:0;top:0;width:960px;height:540px">'
        f"{textboxes}</div>"
    )

    result = DesignQualityEvaluator(get_ruleset()).evaluate([{"index": 1, "html": html}], {})

    messages = [issue["message"] for issue in result.per_slide[0]["issues"]]
    assert any("Typography lacks visual hierarchy" in message for message in messages)
    assert any("Typography is too plain" in message for message in messages)


def test_pptx_textbox_uses_precise_padding_shorthand() -> None:
    element = ParsedElement(
        pptx_type="textbox",
        pptx_shape=None,
        position={"left": 40, "top": 100, "width": 220, "height": 80},
        styles={
            "padding": "6px 18px 10px 24px",
            "font-size": "14px",
            "color": "#112233",
        },
        text_content="Precise padding",
        children=[],
        attributes={},
    )
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])

    CSStoOOXMLEngine()._add_element(slide, element)
    frame = slide.shapes[0].text_frame

    assert abs(frame.margin_left - Emu(24 * 9525)) < 5
    assert abs(frame.margin_right - Emu(18 * 9525)) < 5
    assert abs(frame.margin_top - Emu(6 * 9525)) < 5
    assert abs(frame.margin_bottom - Emu(10 * 9525)) < 5


def test_pptx_textbox_uses_explicit_text_alignment_attrs() -> None:
    element = ParsedElement(
        pptx_type="textbox",
        pptx_shape=None,
        position={"left": 40, "top": 100, "width": 220, "height": 56},
        styles={
            "font-size": "24px",
            "font-weight": "700",
            "color": "#112233",
        },
        text_content="69.5%",
        children=[],
        attributes={
            "data-pptx-text-role": "kpi_value",
            "data-pptx-text-align": "center",
            "data-pptx-text-valign": "middle",
            "data-pptx-text-padding": '{"top":0,"right":4,"bottom":0,"left":4}',
        },
    )
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])

    CSStoOOXMLEngine()._add_element(slide, element)
    frame = slide.shapes[0].text_frame

    assert frame.vertical_anchor == MSO_ANCHOR.MIDDLE
    assert frame.paragraphs[0].alignment == PP_ALIGN.CENTER
    assert abs(frame.margin_left - Emu(4 * 9525)) < 5
    assert abs(frame.margin_top - Emu(0)) < 5


def test_normalized_html_adds_precise_text_attrs_for_metric() -> None:
    html = """<div data-slide="1" style="position:absolute;left:0;top:0;width:960px;height:540px">
    <div data-pptx-type="textbox"
      style="position:absolute;left:80px;top:120px;width:180px;height:56px;font-size:32px;font-weight:700;color:#6366F1">69.5%</div>
    </div>"""

    normalized = _normalize_slide_html(html)

    assert 'data-pptx-text-role="kpi_value"' in normalized
    assert 'data-pptx-text-align="center"' in normalized
    assert 'data-pptx-text-valign="middle"' in normalized
    assert 'data-pptx-text-padding="0px 4px 0px 4px"' in normalized
    assert "text-align:center" in normalized
    assert "vertical-align:middle" in normalized
    assert "padding:0px 4px 0px 4px" in normalized


def test_normalized_html_mirrors_explicit_alignment_attrs_to_preview_css() -> None:
    html = """<div data-slide="1" style="position:absolute;left:0;top:0;width:960px;height:540px">
    <div data-pptx-type="textbox"
      data-pptx-text-role="kpi_value"
      data-pptx-text-align="center"
      data-pptx-text-valign="middle"
      data-pptx-text-padding='{"top":0,"right":6,"bottom":0,"left":6}'
      style="position:absolute;left:80px;top:120px;width:180px;height:56px;font-size:32px;font-weight:700;color:#6366F1;text-align:left">6종</div>
    </div>"""

    normalized = _normalize_slide_html(html)

    assert "text-align:center" in normalized
    assert "vertical-align:middle" in normalized
    assert "padding:0px 6px 0px 6px" in normalized
    assert 'data-pptx-text-align="center"' in normalized
    assert 'data-pptx-text-valign="middle"' in normalized


def test_preview_css_vertically_centers_textboxes_with_alignment_attrs() -> None:
    assert '[data-pptx-type="textbox"][data-pptx-text-valign="middle"]' in PPTX_TEXT_ALIGNMENT_PREVIEW_CSS
    assert "justify-content:center" in PPTX_TEXT_ALIGNMENT_PREVIEW_CSS


def test_text_is_hard_wrapped_to_card_width() -> None:
    engine = CSStoOOXMLEngine()
    size, line_height, lines = engine._fit_text_lines_to_box(
        text="• 클라우드 MSP 전문 기업 토스뱅크 주주 참여 금융 진출",
        font_size_px=14,
        line_height=1.35,
        container_w=190,
        container_h=96,
        pad_left=12,
        pad_right=12,
        pad_top=10,
        pad_bottom=10,
    )

    assert size <= 14
    assert line_height <= 1.35
    assert len(lines) >= 2
    assert all(len(line) <= 18 for line in lines)


def test_pptx_text_fit_protects_short_colon_label_from_line_break() -> None:
    engine = CSStoOOXMLEngine()
    text = engine._protect_label_colon_breaks("핵심: 채널은 3가지로 다양하지만")
    _size, _line_height, lines = engine._fit_text_lines_to_box(
        text=text,
        font_size_px=13,
        line_height=1.35,
        container_w=54,
        container_h=80,
        pad_left=0,
        pad_right=0,
        pad_top=0,
        pad_bottom=0,
    )

    assert lines[0] != "핵심:"
    assert lines[0].startswith("핵심:\u00a0")


def test_pptx_text_fit_preserves_natural_single_paragraph() -> None:
    engine = CSStoOOXMLEngine()
    text = "Korean and English mixed content should fit by shrinking instead of forced paragraph breaks"
    size, line_height, lines = engine._fit_text_lines_to_box(
        text=text,
        font_size_px=14,
        line_height=1.35,
        container_w=180,
        container_h=72,
        pad_left=10,
        pad_right=10,
        pad_top=8,
        pad_bottom=8,
        preserve_natural_lines=True,
    )

    assert size <= 14
    assert line_height <= 1.35
    assert lines == [text]


def test_pptx_single_paragraph_text_is_not_forced_into_multiple_paragraphs() -> None:
    element = ParsedElement(
        pptx_type="textbox",
        pptx_shape=None,
        position={"left": 40, "top": 100, "width": 140, "height": 64},
        styles={"font-size": "14px", "line-height": "1.35", "color": "#111827"},
        text_content="HTML과 PPTX의 상하 자간 차이를 줄입니다",
        children=[],
        attributes={},
    )
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])

    CSStoOOXMLEngine()._add_element(slide, element)
    paragraphs = slide.shapes[0].text_frame.paragraphs

    assert len(paragraphs) == 1
    assert all(paragraph.space_before is None or paragraph.space_before.pt == 0 for paragraph in paragraphs)
    assert all(paragraph.space_after is None or paragraph.space_after.pt == 0 for paragraph in paragraphs)


def test_chart_table_shape_options_do_not_break_native_conversion(tmp_path: Path) -> None:
    html = """<div data-slide="1" style="position:absolute;left:0;top:0;width:960px;height:540px">
    <div data-pptx-type="chart" data-pptx-chart-type="bar"
      data-pptx-chart-data='[{"label":"2022","value":"2,600"},{"label":"2023","value":"3,350"}]'
      data-pptx-chart-options='{"show_legend":false,"data_label_position":"outside_end","gap_width":40,"colors":["#3B82F6"]}'
      style="position:absolute;left:40px;top:80px;width:360px;height:220px"></div>
    <div data-pptx-type="table"
      data-pptx-table-data='{"headers":["항목","값"],"rows":[["매출","4,058"],["성장률","21%"]]}'
      data-pptx-table-options='{"numeric_align":"right","vertical_align":"middle","body_font_size":9}'
      style="position:absolute;left:430px;top:80px;width:300px;height:120px"></div>
    <div data-pptx-type="shape" data-pptx-shape="diamond"
      data-pptx-shape-options='{"line_color":"#FF0000","line_width":2,"transparency":0.2}'
      style="position:absolute;left:760px;top:100px;width:40px;height:40px;background-color:#10B981"></div>
    </div>"""

    output = CSStoOOXMLEngine().build_presentation(
        [{"index": 1, "html": html, "metadata": {"slide_type": "content"}}],
        tmp_path,
    )

    assert output.exists()
    assert output.stat().st_size > 0


def test_line_like_shapes_do_not_receive_pptx_shadow() -> None:
    element = ParsedElement(
        pptx_type="shape",
        pptx_shape="rect",
        position={"left": 40, "top": 120, "width": 360, "height": 2},
        styles={
            "background-color": "#10B981",
            "box-shadow": "0px 2px 6px rgba(0,0,0,0.25)",
        },
        text_content="",
        children=[],
        attributes={},
    )
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])

    CSStoOOXMLEngine()._add_element(slide, element)

    assert "outerShdw" not in slide.shapes[0]._element.xml


def test_arrow_shapes_receive_stable_geometry_and_no_shadow() -> None:
    element = ParsedElement(
        pptx_type="shape",
        pptx_shape="right_arrow",
        position={"left": 40, "top": 120, "width": 6, "height": 4},
        styles={
            "color": "#10B981",
            "box-shadow": "0px 2px 6px rgba(0,0,0,0.25)",
        },
        text_content="",
        children=[],
        attributes={},
    )
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])

    CSStoOOXMLEngine()._add_element(slide, element)

    assert element.position["width"] == 18.0
    assert element.position["height"] == 10.0
    assert "outerShdw" not in slide.shapes[0]._element.xml
    assert slide.shapes[0].shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE


def test_fallback_icon_is_created_when_iconify_cache_is_missing() -> None:
    path = get_fallback_icon_path("database", color="112233", size=32)

    assert path is not None
    assert path.exists()
    with Image.open(path) as image:
        assert image.mode == "RGBA"
        assert image.getpixel((0, 0))[3] == 0


def test_architecture_request_routes_to_diagrams_visual_asset() -> None:
    plan = _fallback_plan(
        "AWS architecture 3-tier diagram을 다시 그려줘",
        [{"index": 1, "slide_type": "cover"}, {"index": 2, "slide_type": "content"}],
    )

    assert plan["enabled"] is True
    assert plan["assets"][0]["slide_index"] == 2
    assert plan["assets"][0]["method"] == METHOD_DIAGRAMS
    assert plan["assets"][0]["diagrams_provider"] == "aws"
    assert plan["assets"][0]["diagrams_nodes"]
    assert "CloudFront" in plan["assets"][0]["mermaid"]


def test_legacy_mermaid_without_structured_topology_is_rejected() -> None:
    plan = _normalize_plan(
        {
            "enabled": True,
            "assets": [
                {
                    "slide_index": 2,
                    "method": "mermaid_image",
                    "description": "Logical workflow diagram",
                    "mermaid": "graph LR\n  A[Input] --> B[Output]",
                }
            ],
        },
        "workflow diagram",
        [{"index": 2, "slide_type": "content"}],
    )

    assert plan["enabled"] is False
    assert plan["assets"] == []


def test_visual_asset_placement_uses_planned_image_slot() -> None:
    plan = _normalize_plan(
        {
            "enabled": True,
            "assets": [
                {
                    "slide_index": 2,
                    "method": METHOD_DIAGRAMS,
                    "description": "LangGraph architecture",
                    "diagrams_nodes": [
                        {"id": "user", "label": "User", "provider": "generic", "service": "users"},
                        {"id": "app", "label": "LangGraph App", "provider": "generic", "service": "agent"},
                    ],
                    "diagrams_edges": [{"from": "user", "to": "app", "label": "Request"}],
                    "placement": {"x": 360, "y": 112, "w": 300, "h": 180},
                }
            ],
        },
        "LangGraph architecture diagram",
        [
            {
                "index": 2,
                "slide_type": "content",
                "layout_plan": {
                    "element_placements": [
                        {
                            "id": "visual_asset_main",
                            "element": "image",
                            "asset_role": "visual_asset",
                            "x": 48,
                            "y": 92,
                            "w": 604,
                            "h": 318,
                        }
                    ]
                },
            }
        ],
    )

    assert plan["assets"][0]["placement"] == {"x": 48, "y": 92, "w": 604, "h": 318}


def test_visual_asset_slot_constraints_expose_aspect_guidance() -> None:
    tall = _slot_constraints({"x": 120, "y": 90, "w": 180, "h": 420})
    wide = _slot_constraints({"x": 48, "y": 92, "w": 604, "h": 318})

    assert tall["recommended_direction"] == "TB"
    assert tall["max_recommended_nodes"] <= 5
    assert wide["recommended_direction"] == "LR"
    assert wide["slot_aspect_ratio"] > 1.4


def test_visual_asset_direction_is_adjusted_to_tall_slot() -> None:
    plan = _normalize_plan(
        {
            "enabled": True,
            "assets": [
                {
                    "slide_index": 2,
                    "method": METHOD_DIAGRAMS,
                    "description": "Tall architecture",
                    "diagrams_direction": "LR",
                    "diagrams_nodes": [
                        {"id": "a", "label": "A", "provider": "generic", "service": "process"},
                        {"id": "b", "label": "B", "provider": "generic", "service": "process"},
                        {"id": "c", "label": "C", "provider": "generic", "service": "process"},
                    ],
                    "diagrams_edges": [
                        {"from": "a", "to": "b"},
                        {"from": "b", "to": "c"},
                    ],
                    "placement": {"x": 120, "y": 90, "w": 180, "h": 420},
                }
            ],
        },
        "tall diagram",
        [{"index": 2, "slide_type": "content"}],
    )

    assert plan["assets"][0]["diagrams_direction"] == "TB"
    assert plan["assets"][0]["render_profile"]["recommended_direction"] == "TB"


def test_diagram_png_canvas_matches_slot_aspect(tmp_path: Path) -> None:
    path = tmp_path / "diagram.png"
    Image.new("RGBA", (400, 120), (12, 34, 56, 255)).save(path)

    _fit_png_canvas_to_placement(path, {"x": 0, "y": 0, "w": 180, "h": 420})

    with Image.open(path) as image:
        width, height = image.size
    assert abs((width / height) - (180 / 420)) < 0.03


def test_diagram_png_border_is_baked_into_image(tmp_path: Path) -> None:
    path = tmp_path / "diagram.png"
    Image.new("RGBA", (900, 520), (255, 255, 255, 0)).save(path)

    _add_diagram_image_border(path)

    with Image.open(path) as image:
        pixel = image.convert("RGBA").getpixel((0, 0))
    assert pixel[3] > 0


def test_diagram_rendering_uses_high_resolution_defaults() -> None:
    assert MIN_DIAGRAM_PNG_WIDTH >= 3600
    assert MIN_DIAGRAM_PNG_HEIGHT >= 2250
    assert DIAGRAMS_DPI >= 450


def test_llm_plan_validation_detects_missing_reserved_slot_asset() -> None:
    issues = _llm_plan_validation_issues(
        {
            "enabled": True,
            "assets": [
                {
                    "slide_index": 2,
                    "method": METHOD_DIAGRAMS,
                    "description": "Architecture",
                    "diagrams_nodes": [
                        {"id": "app", "label": "App", "provider": "generic", "service": "agent"},
                        {"id": "db", "label": "DB", "provider": "generic", "service": "database"},
                    ],
                    "diagrams_edges": [{"from": "app", "to": "db", "label": "SQL"}],
                }
            ],
        },
        "Create diagrams for reserved slots",
        [
            {
                "index": 2,
                "layout_plan": {
                    "element_placements": [
                        {"element": "image", "asset_role": "visual_asset", "x": 40, "y": 90, "w": 400, "h": 240}
                    ]
                },
            },
            {
                "index": 4,
                "title": "Architecture diagram",
                "suggested_elements": ["diagram"],
                "layout_plan": {
                    "element_placements": [
                        {"element": "image", "asset_role": "visual_asset", "x": 40, "y": 90, "w": 400, "h": 240}
                    ]
                },
            },
        ],
    )

    assert any("slides: 4" in issue for issue in issues)


def test_llm_plan_validation_allows_empty_assets_for_soft_visual_signal() -> None:
    issues = _llm_plan_validation_issues(
        {"enabled": False, "assets": []},
        "Create a product brief with a workflow overview",
        [
            {
                "index": 2,
                "title": "Generation workflow",
                "key_message": "The workflow can be represented with native cards and connectors.",
                "suggested_elements": ["connector", "card"],
                "layout_plan": {
                    "element_placements": [
                        {
                            "element": "image",
                            "asset_role": "visual_asset",
                            "x": 48,
                            "y": 92,
                            "w": 604,
                            "h": 318,
                        }
                    ]
                },
            }
        ],
    )

    assert issues == []


def test_llm_plan_validation_requires_assets_for_explicit_redraw_request() -> None:
    issues = _llm_plan_validation_issues(
        {"enabled": False, "assets": []},
        "첨부된 아키텍처 다이어그램을 다시 그려줘",
        [{"index": 2, "title": "Architecture diagram", "suggested_elements": ["diagram"]}],
    )

    assert any("explicit visual asset request" in issue for issue in issues)


def test_workflow_topic_alone_does_not_trigger_external_visual_asset() -> None:
    assert _quick_visual_signal(
        "Create a product brief",
        [
            {
                "index": 2,
                "title": "Agentic workflow",
                "key_message": "The process is shown with native cards.",
                "suggested_elements": ["connector", "card"],
            }
        ],
    ) is False


def test_llm_plan_validation_detects_duplicate_diagram_topologies() -> None:
    asset = {
        "method": METHOD_DIAGRAMS,
        "description": "Repeated topology",
        "diagrams_nodes": [
            {"id": "client", "label": "Client", "provider": "generic", "service": "users"},
            {"id": "app", "label": "Application Service", "provider": "generic", "service": "agent"},
            {"id": "data", "label": "Data Store", "provider": "generic", "service": "database"},
        ],
        "diagrams_edges": [
            {"from": "client", "to": "app", "label": "Request"},
            {"from": "app", "to": "data", "label": "Data"},
        ],
    }
    issues = _llm_plan_validation_issues(
        {
            "enabled": True,
            "assets": [
                {"slide_index": 3, **asset},
                {"slide_index": 5, **asset},
            ],
        },
        "Create slide-specific diagrams",
        [{"index": 3}, {"index": 5}],
    )

    assert any("Duplicate diagrams topology" in issue for issue in issues)


def test_truncated_json_detector_flags_unbalanced_output() -> None:
    assert _looks_like_truncated_json('{"enabled": true, "assets": [{"slide_index": 2}')
    assert not _looks_like_truncated_json('{"enabled": true, "assets": []}')


def test_visual_asset_slot_fallback_overrides_llm_skip() -> None:
    plan = _fallback_visual_slot_plan(
        "Create a product brief",
        [
            {
                "index": 2,
                "title": "Service architecture",
                "key_message": "Architecture diagram is required to explain service topology.",
                "suggested_elements": ["diagram"],
                "layout_plan": {
                    "element_placements": [
                        {
                            "id": "visual_asset_main",
                            "element": "image",
                            "asset_role": "visual_asset",
                            "x": 48,
                            "y": 92,
                            "w": 604,
                            "h": 318,
                        }
                    ]
                },
            }
        ],
    )

    assert plan["enabled"] is True
    assert plan["assets"][0]["method"] == METHOD_DIAGRAMS
    assert plan["assets"][0]["placement"]["w"] == 604


def test_fallback_visual_slots_create_slide_specific_diagrams() -> None:
    blueprints = [
        {
            "index": 3,
            "title": "Native document conversion workflow diagram",
            "key_message": "A flowchart is required to explain how natural language becomes native document output.",
            "suggested_elements": ["flowchart"],
            "layout_plan": {
                "element_placements": [
                    {
                        "id": "visual_asset_main",
                        "element": "image",
                        "asset_role": "visual_asset",
                        "x": 80,
                        "y": 120,
                        "w": 520,
                        "h": 280,
                    }
                ]
            },
        },
        {
            "index": 5,
            "title": "PPTX vs Rich Document agentic pipeline diagram",
            "key_message": "A pipeline diagram is required to compare PPTX and rich document generation.",
            "suggested_elements": ["diagram"],
            "layout_plan": {
                "element_placements": [
                    {
                        "id": "visual_asset_main",
                        "element": "image",
                        "asset_role": "visual_asset",
                        "x": 80,
                        "y": 120,
                        "w": 520,
                        "h": 280,
                    }
                ]
            },
        },
    ]

    plan = _fallback_visual_slot_plan("DocuMind product brief", blueprints)
    labels_by_slide = {
        asset["slide_index"]: [node["label"] for node in asset["diagrams_nodes"]]
        for asset in plan["assets"]
    }

    assert labels_by_slide[3] != labels_by_slide[5]
    assert "User Request" in labels_by_slide[3]
    assert "PPTX Pipeline" in labels_by_slide[5]
    assert "Native Document Output" in labels_by_slide[5]


def test_generic_fallback_diagram_edges_follow_content_nodes() -> None:
    text = "Web UI and REST API feed a LangGraph agent, QA validation, and document output"

    nodes = _fallback_diagrams_nodes(text, METHOD_DIAGRAMS)
    edges = _fallback_diagrams_edges(text, METHOD_DIAGRAMS)

    assert [node["label"] for node in nodes] != [
        "Client",
        "Edge / Gateway",
        "Application Service",
        "Data Store",
    ]
    assert len(edges) == len(nodes) - 1
    assert edges[0]["from"] == nodes[0]["id"]
    assert edges[-1]["to"] == nodes[-1]["id"]


def test_visual_asset_signal_ignores_generic_reserved_image_slot() -> None:
    assert _quick_visual_signal(
        "Create a product brief",
        [
            {
                "index": 2,
                "layout_plan": {
                    "element_placements": [
                        {
                            "id": "visual_asset_main",
                            "element": "image",
                            "asset_role": "visual_asset",
                            "x": 48,
                            "y": 92,
                            "w": 604,
                            "h": 318,
                        }
                    ]
                },
            }
        ],
    ) is False


def test_visual_asset_signal_uses_architecture_context_not_slot_alone() -> None:
    blueprint = {
        "index": 2,
        "title": "Service architecture diagram",
        "key_message": "Diagram is required to explain service topology.",
        "suggested_elements": ["diagram"],
        "layout_plan": {
            "element_placements": [
                {
                    "id": "visual_asset_main",
                    "element": "image",
                    "asset_role": "visual_asset",
                    "x": 48,
                    "y": 92,
                    "w": 604,
                    "h": 318,
                }
            ]
        },
    }

    assert _slide_context_requires_diagram(blueprint, "Create a product brief") is True
    assert _quick_visual_signal("Create a product brief", [blueprint]) is True


def test_visual_asset_plan_adds_missing_reserved_slot_asset() -> None:
    blueprints = [
        {
            "index": 2,
            "title": "Architecture diagram",
            "key_message": "Diagram is required to explain the architecture.",
            "suggested_elements": ["diagram"],
            "layout_plan": {
                "element_placements": [
                    {
                        "id": "visual_asset_main",
                        "element": "image",
                        "asset_role": "visual_asset",
                        "x": 48,
                        "y": 92,
                        "w": 400,
                        "h": 260,
                    }
                ]
            },
        },
        {
            "index": 4,
            "title": "Operations flow diagram",
            "key_message": "Flowchart is required to explain operations.",
            "suggested_elements": ["flowchart"],
            "layout_plan": {
                "element_placements": [
                    {
                        "id": "visual_asset_main",
                        "element": "image",
                        "asset_role": "visual_asset",
                        "x": 64,
                        "y": 130,
                        "w": 560,
                        "h": 300,
                    }
                ]
            },
        },
    ]
    normalized = {
        "enabled": True,
        "reason": "LLM selected one asset",
        "assets": [
            {
                "id": "asset_2_1",
                "slide_index": 2,
                "method": METHOD_DIAGRAMS,
                "description": "Architecture diagram",
            }
        ],
    }

    merged = _merge_missing_reserved_slot_assets(normalized, "Create a product brief", blueprints)

    assert [asset["slide_index"] for asset in merged["assets"]] == [2, 4]
    assert merged["assets"][1]["placement"] == {"x": 64, "y": 130, "w": 560, "h": 300}


def test_diagrams_topology_fields_are_preserved() -> None:
    plan = _normalize_plan(
        {
            "enabled": True,
            "assets": [
                {
                    "slide_index": 2,
                    "method": METHOD_DIAGRAMS,
                    "description": "Azure application architecture",
                    "diagrams_provider": "azure",
                    "diagrams_direction": "TB",
                    "diagrams_clusters": [
                        {"id": "region", "label": "Azure Region"},
                        {"id": "vnet", "label": "Virtual Network", "parent": "region"},
                    ],
                    "diagrams_nodes": [
                        {
                            "id": "APP",
                            "library": "azure",
                            "service": "app_services",
                            "label": "App Service",
                            "cluster": "vnet",
                        },
                        {
                            "id": "DB",
                            "provider": "azure",
                            "service": "sql_database",
                            "label": "SQL Database",
                            "cluster": "vnet",
                        },
                    ],
                    "diagrams_edges": [
                        {
                            "from": "APP",
                            "to": "DB",
                            "label": "SQL",
                            "style": "dashed",
                        }
                    ],
                    "mermaid": "graph LR\n  U[Users] --> APP[App Service]\n  APP --> DB[Cloud SQL]",
                }
            ],
        },
        "Azure 기반으로 구성해줘",
        [{"index": 2, "slide_type": "content"}],
    )

    asset = plan["assets"][0]
    topology = _diagrams_topology(asset)
    assert asset["diagrams_provider"] == "azure"
    assert topology["direction"] == "TB"
    assert topology["clusters"][1] == {
        "id": "vnet",
        "label": "Virtual Network",
        "parent": "region",
    }
    assert topology["nodes"][0]["provider"] == "azure"
    assert topology["nodes"][0]["service"] == "appservice"
    assert topology["nodes"][1]["service"] == "sqldatabase"
    assert topology["edges"][0]["style"] == "dashed"


def test_diagrams_node_candidates_fall_back_to_specific_provider_icons() -> None:
    candidates = _diagrams_node_candidates(
        {"provider": "generic", "service": "appservice", "label": "Azure App Service"}
    )

    assert ("diagrams.azure.compute", "AppServices") in candidates
    assert candidates[-1] == ("diagrams.generic.blank", "Blank")


def test_generic_diagrams_prefers_flowchart_icons_over_blank() -> None:
    agent_candidates = _diagrams_node_candidates(
        {"provider": "generic", "service": "agent", "label": "LangChain Agent"}
    )
    docs_candidates = _diagrams_node_candidates(
        {"provider": "generic", "service": "document", "label": "Retrieved Docs"}
    )
    search_candidates = _diagrams_node_candidates(
        {"provider": "generic", "service": "search", "label": "Retrieval & Search"}
    )

    assert ("diagrams.programming.flowchart", "PredefinedProcess") in agent_candidates
    assert ("diagrams.programming.flowchart", "Document") in docs_candidates
    assert ("diagrams.programming.flowchart", "Inspection") in search_candidates
    assert agent_candidates.index(("diagrams.programming.flowchart", "PredefinedProcess")) < agent_candidates.index(("diagrams.generic.blank", "Blank"))


def test_visual_asset_plan_honors_targeted_slide_instruction() -> None:
    plan = _normalize_plan(
        {
            "enabled": True,
            "assets": [
                {
                    "slide_index": 5,
                    "method": METHOD_DIAGRAMS,
                    "title": "LangChain RAG Pipeline Architecture",
                    "description": "LangChain RAG Pipeline architecture diagram",
                    "diagrams_nodes": [
                        {"id": "query", "label": "Query", "provider": "generic", "service": "input"},
                        {"id": "rag", "label": "RAG Pipeline", "provider": "generic", "service": "agent"},
                    ],
                    "diagrams_edges": [{"from": "query", "to": "rag", "label": "Search"}],
                }
            ],
        },
        "슬라이드 3: LangChain RAG Pipeline 아키텍처 다이어그램으로 변경해줘. "
        "슬라이드 5: 하단 카드 영역 구분.",
        [{"index": 3, "slide_type": "content"}, {"index": 5, "slide_type": "content"}],
        {
            3: "LangChain RAG Pipeline 아키텍처 다이어그램으로 변경해줘.",
            5: "하단 카드 영역 구분.",
        },
    )

    assert plan["assets"][0]["slide_index"] == 3


def test_safe_diagrams_topology_breaks_cluster_cycles_and_bad_refs() -> None:
    topology = _safe_diagrams_topology(
        {
            "method": METHOD_DIAGRAMS,
            "title": "Cyclic Architecture",
            "diagrams_clusters": [
                {"id": "a", "label": "A", "parent": "b"},
                {"id": "b", "label": "B", "parent": "a"},
                {"id": "c", "label": "C", "parent": "missing"},
            ],
            "diagrams_nodes": [
                {"id": "app", "label": "App", "provider": "aws", "service": "ec2", "cluster": "a"},
                {"id": "db", "label": "DB", "provider": "aws", "service": "rds", "cluster": "missing"},
            ],
            "diagrams_edges": [
                {"from": "app", "to": "db", "label": "SQL"},
                {"from": "db", "to": "db", "label": "self"},
                {"from": "app", "to": "missing", "label": "bad"},
            ],
        }
    )

    parent_by_id = {cluster["id"]: cluster["parent"] for cluster in topology["clusters"]}
    for cluster_id in parent_by_id:
        seen = {cluster_id}
        parent = parent_by_id[cluster_id]
        while parent:
            assert parent not in seen
            seen.add(parent)
            parent = parent_by_id.get(parent, "")

    assert {edge["label"] for edge in topology["edges"]} == {"SQL"}
    assert next(node for node in topology["nodes"] if node["id"] == "db")["cluster"] == ""


async def test_diagrams_asset_does_not_fallback_when_native_renderer_recurses(
    tmp_path: Path,
    monkeypatch,
) -> None:
    def boom(asset, output_path):
        raise RecursionError("maximum recursion depth exceeded")

    monkeypatch.setattr(
        "src.formats.pptx.agents.nodes.visual_asset_planner._render_diagrams_asset",
        boom,
    )
    asset = {
        "id": "recursing",
        "slide_index": 2,
        "method": METHOD_DIAGRAMS,
        "title": "Safe fallback",
        "description": "AWS architecture",
        "diagrams_nodes": [
            {"id": "app", "label": "App", "provider": "aws", "service": "ec2"},
            {"id": "db", "label": "DB", "provider": "aws", "service": "rds"},
        ],
        "diagrams_edges": [{"from": "app", "to": "db", "label": "SQL"}],
    }

    rendered = await _render_asset(asset, tmp_path)

    assert rendered is None
    assert not any(tmp_path.glob("recursing_*.png"))


def test_db_backed_pptx_preview_embeds_local_visual_asset(tmp_path: Path) -> None:
    image_path = tmp_path / "diagram.png"
    Image.new("RGB", (64, 48), (12, 34, 56)).save(image_path)
    html = (
        '<div data-slide="1">'
        f'<div data-pptx-type="image" data-pptx-image-path="{image_path}" '
        'style="position:absolute;left:40px;top:100px;width:300px;height:180px"></div>'
        '</div>'
    )

    preview = _embed_pptx_preview_assets(html)

    assert "data:image/png;base64" in preview
    assert str(image_path) not in preview


def test_visual_asset_injection_skips_missing_slot_to_avoid_overlay(tmp_path: Path) -> None:
    image_path = tmp_path / "diagram.png"
    Image.new("RGB", (64, 48), (12, 34, 56)).save(image_path)
    html = (
        '<div data-slide="2">'
        '<div data-pptx-region="header" data-pptx-type="textbox" '
        'style="position:absolute;left:40px;top:20px;width:500px;height:30px">Title</div>'
        '<div data-pptx-type="textbox" style="position:absolute;left:120px;top:130px;'
        'width:240px;height:80px">Overlapping content</div>'
        '<div data-pptx-type="textbox" style="position:absolute;left:40px;top:440px;'
        'width:200px;height:40px">Safe note</div>'
        "</div>"
    )
    slides = [{"index": 2, "html": html, "elements_used": []}]
    assets = [
        {
            "id": "arch",
            "slide_index": 2,
            "path": str(image_path),
            "placement": {"x": 100, "y": 110, "w": 320, "h": 180},
        }
    ]

    [slide] = _inject_visual_asset_images(slides, assets)

    assert 'data-pptx-image-id="arch"' not in slide["html"]
    assert "Overlapping content" in slide["html"]
    assert "Safe note" in slide["html"]
    assert "Title" in slide["html"]


def test_visual_asset_injection_fills_existing_slot(tmp_path: Path) -> None:
    image_path = tmp_path / "diagram.png"
    Image.new("RGB", (64, 48), (12, 34, 56)).save(image_path)
    html = (
        '<div data-slide="2">'
        '<div data-pptx-type="shape" data-pptx-shape="rounded_rect" '
        'data-pptx-asset-id="arch" '
        'style="position:absolute;left:120px;top:130px;width:320px;height:180px"></div>'
        '<div data-pptx-type="textbox" style="position:absolute;left:40px;top:440px;'
        'width:200px;height:40px">Safe note</div>'
        "</div>"
    )
    slides = [{"index": 2, "html": html, "elements_used": []}]
    assets = [
        {
            "id": "arch",
            "slide_index": 2,
            "path": str(image_path),
            "placement": {"x": 100, "y": 110, "w": 320, "h": 180},
        }
    ]

    [slide] = _inject_visual_asset_images(slides, assets)

    assert 'data-pptx-type="image"' in slide["html"]
    assert f'data-pptx-image-path="{image_path}"' in slide["html"]
    assert 'data-pptx-image-fit="contain"' in slide["html"]
    assert slide["html"].count('data-pptx-image-id="arch"') == 1
    assert "Safe note" in slide["html"]


def test_visual_asset_injection_fills_geometry_matched_empty_slot(tmp_path: Path) -> None:
    image_path = tmp_path / "diagram.png"
    Image.new("RGB", (64, 48), (12, 34, 56)).save(image_path)
    html = (
        '<div data-slide="2">'
        '<div data-pptx-type="shape" data-pptx-shape="rounded_rect" '
        'style="position:absolute;left:100px;top:110px;width:320px;height:180px;'
        'background-color:#E0F2FE;border:1px solid #93C5FD"></div>'
        '<div data-pptx-type="textbox" style="position:absolute;left:40px;top:440px;'
        'width:200px;height:40px">Safe note</div>'
        "</div>"
    )
    slides = [{"index": 2, "html": html, "elements_used": []}]
    assets = [
        {
            "id": "arch",
            "slide_index": 2,
            "path": str(image_path),
            "placement": {"x": 100, "y": 110, "w": 320, "h": 180},
        }
    ]

    [slide] = _inject_visual_asset_images(slides, assets)

    assert 'data-pptx-type="image"' in slide["html"]
    assert f'data-pptx-image-path="{image_path}"' in slide["html"]
    assert slide["html"].count('data-pptx-image-id="arch"') == 1
    assert "Safe note" in slide["html"]


def test_visual_asset_injection_does_not_replace_populated_geometry_slot(tmp_path: Path) -> None:
    image_path = tmp_path / "diagram.png"
    Image.new("RGB", (64, 48), (12, 34, 56)).save(image_path)
    html = (
        '<div data-slide="2">'
        '<div data-pptx-type="shape" data-pptx-shape="rounded_rect" '
        'style="position:absolute;left:100px;top:110px;width:320px;height:180px;'
        'background-color:#E0F2FE;border:1px solid #93C5FD"></div>'
        '<div data-pptx-type="textbox" style="position:absolute;left:128px;top:140px;'
        'width:220px;height:40px">Existing content</div>'
        "</div>"
    )
    slides = [{"index": 2, "html": html, "elements_used": []}]
    assets = [
        {
            "id": "arch",
            "slide_index": 2,
            "path": str(image_path),
            "placement": {"x": 100, "y": 110, "w": 320, "h": 180},
        }
    ]

    [slide] = _inject_visual_asset_images(slides, assets)

    assert 'data-pptx-image-id="arch"' not in slide["html"]
    assert "Existing content" in slide["html"]


def test_visual_asset_injection_dedupes_duplicate_image_nodes(tmp_path: Path) -> None:
    image_path = tmp_path / "diagram.png"
    Image.new("RGB", (64, 48), (12, 34, 56)).save(image_path)
    html = (
        '<div data-slide="2">'
        f'<div data-pptx-type="image" data-pptx-image-id="arch" data-pptx-image-path="{image_path}" '
        'style="position:absolute;left:40px;top:100px;width:300px;height:180px"></div>'
        f'<div data-pptx-type="image" data-pptx-image-id="arch" data-pptx-image-path="{image_path}" '
        'style="position:absolute;left:420px;top:100px;width:300px;height:180px"></div>'
        "</div>"
    )
    slides = [{"index": 2, "html": html, "elements_used": []}]
    assets = [{"id": "arch", "slide_index": 2, "path": str(image_path)}]

    [slide] = _inject_visual_asset_images(slides, assets)

    assert slide["html"].count('data-pptx-image-id="arch"') == 1


def test_visual_asset_injection_does_not_overlay_manual_diagram(tmp_path: Path) -> None:
    image_path = tmp_path / "diagram.png"
    Image.new("RGB", (64, 48), (12, 34, 56)).save(image_path)
    html = (
        '<div data-slide="2">'
        '<div data-pptx-type="shape" data-pptx-shape="rounded_rect" '
        'style="position:absolute;left:40px;top:100px;width:120px;height:60px"></div>'
        '<div data-pptx-type="shape" data-pptx-shape="rounded_rect" '
        'style="position:absolute;left:240px;top:100px;width:120px;height:60px"></div>'
        '<div data-pptx-type="connector" style="position:absolute;left:160px;top:128px;width:80px;height:2px"></div>'
        '<div data-pptx-type="connector" style="position:absolute;left:360px;top:128px;width:80px;height:2px"></div>'
        "</div>"
    )
    slides = [{"index": 2, "html": html, "elements_used": []}]
    assets = [{"id": "arch", "slide_index": 2, "path": str(image_path)}]

    [slide] = _inject_visual_asset_images(slides, assets)

    assert 'data-pptx-image-id="arch"' not in slide["html"]


def test_visual_asset_injection_fills_slot_even_with_manual_diagram(tmp_path: Path) -> None:
    image_path = tmp_path / "diagram.png"
    Image.new("RGB", (64, 48), (12, 34, 56)).save(image_path)
    html = (
        '<div data-slide="2">'
        '<div data-pptx-type="shape" data-pptx-shape="rounded_rect" '
        'data-pptx-asset-role="visual_asset" '
        'style="position:absolute;left:120px;top:130px;width:320px;height:180px"></div>'
        '<div data-pptx-type="shape" data-pptx-shape="rounded_rect" '
        'style="position:absolute;left:40px;top:100px;width:120px;height:60px"></div>'
        '<div data-pptx-type="shape" data-pptx-shape="rounded_rect" '
        'style="position:absolute;left:240px;top:100px;width:120px;height:60px"></div>'
        '<div data-pptx-type="connector" style="position:absolute;left:160px;top:128px;width:80px;height:2px"></div>'
        '<div data-pptx-type="connector" style="position:absolute;left:360px;top:128px;width:80px;height:2px"></div>'
        "</div>"
    )
    slides = [{"index": 2, "html": html, "elements_used": []}]
    assets = [{"id": "arch", "slide_index": 2, "path": str(image_path)}]

    [slide] = _inject_visual_asset_images(slides, assets)

    assert 'data-pptx-type="image"' in slide["html"]
    assert 'data-pptx-image-id="arch"' in slide["html"]


async def test_render_asset_reports_cache_hit_renderer(tmp_path: Path) -> None:
    asset = {
        "id": "asset_cache",
        "slide_index": 2,
        "asset_type": "architecture",
        "method": METHOD_DIAGRAMS,
        "title": "Cached architecture",
        "description": "Cached architecture diagram",
        "diagrams_nodes": [
            {"id": "app", "label": "App", "provider": "generic", "service": "agent"},
            {"id": "db", "label": "DB", "provider": "generic", "service": "database"},
        ],
        "diagrams_edges": [{"from": "app", "to": "db", "label": "SQL"}],
        "placement": {"x": 48, "y": 92, "w": 604, "h": 318},
    }
    import hashlib
    import json

    fingerprint = hashlib.md5(
        json.dumps(asset, ensure_ascii=False, sort_keys=True).encode("utf-8")
    ).hexdigest()[:10]
    cached_path = tmp_path / f"asset_cache_{fingerprint}.png"
    Image.new("RGB", (1800, 1100), (255, 255, 255)).save(cached_path)

    rendered = await _render_asset(asset, tmp_path)

    assert rendered is not None
    assert rendered["renderer"] == "cache_hit"
    assert rendered["cache_hit"] is True


def test_pptx_image_preserves_source_aspect_inside_html_box(tmp_path: Path) -> None:
    image_path = tmp_path / "wide_diagram.png"
    Image.new("RGB", (400, 100), (12, 34, 56)).save(image_path)
    html = (
        '<div data-slide="1">'
        f'<div data-pptx-type="image" data-pptx-image-path="{image_path}" '
        'style="position:absolute;left:40px;top:100px;width:300px;height:300px"></div>'
        '</div>'
    )

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    for element in parse_slide_html(html):
        CSStoOOXMLEngine()._add_element(slide, element)

    picture = slide.shapes[0]
    assert picture.shape_type == MSO_SHAPE_TYPE.PICTURE
    assert abs(picture.left - Emu(40 * 9525)) < 5
    assert abs(picture.top - Emu(212.5 * 9525)) < 5
    assert abs(picture.width - Emu(300 * 9525)) < 5
    assert abs(picture.height - Emu(75 * 9525)) < 5


def test_slide_html_normalizer_shrinks_overflowing_text_and_resolves_overlap() -> None:
    html = (
        '<div data-slide="1" style="position:absolute;left:0;top:0;width:960px;height:540px">'
        '<div data-pptx-type="textbox" style="position:absolute;left:40px;top:90px;'
        'width:360px;height:80px;font-size:54px;line-height:1.1;color:#111827">'
        'AWS 3-Tier 인프라 아키텍처</div>'
        '<div data-pptx-type="textbox" style="position:absolute;left:60px;top:110px;'
        'width:280px;height:60px;font-size:18px;color:#111827">중첩되는 설명 텍스트</div>'
        '</div>'
    )

    output = _normalize_slide_html(html)

    assert "font-size:54px" not in output
    assert "overflow:hidden" in output
    assert "top:176px" in output


def test_slide_html_normalizer_moves_label_off_table_edge_overlap() -> None:
    html = (
        '<div data-slide="1" style="position:absolute;left:0;top:0;width:960px;height:540px">'
        '<div data-pptx-type="table" data-pptx-table-data=\'{"headers":["Phase","설명"],"rows":[["A","B"]]}\' '
        'style="position:absolute;left:40px;top:120px;width:880px;height:180px"></div>'
        '<div data-pptx-type="textbox" style="position:absolute;left:40px;top:286px;'
        'width:360px;height:28px;font-size:15px;font-weight:700;color:#1E1B4B">'
        'Rich Document Agentic Flow</div>'
        '<div data-pptx-type="table" data-pptx-table-data=\'{"headers":["단계","설명"],"rows":[["1","2"]]}\' '
        'style="position:absolute;left:40px;top:304px;width:880px;height:180px"></div>'
        '</div>'
    )

    elements = parse_slide_html(_normalize_slide_html(html))
    first_table, label, second_table = elements[0], elements[1], elements[2]

    assert label.position["top"] >= first_table.position["top"] + first_table.position["height"]
    assert second_table.position["top"] >= label.position["top"] + label.position["height"]


def test_slide_html_normalizer_drops_large_icon_only_empty_card() -> None:
    html = (
        '<div data-slide="1" style="position:absolute;left:0;top:0;width:960px;height:540px">'
        '<div data-pptx-type="shape" data-pptx-shape="rounded_rect" '
        'style="position:absolute;left:80px;top:120px;width:360px;height:150px;'
        'background-color:#FFFFFF;border:1px solid #E5E7EB"></div>'
        '<div data-pptx-type="icon" data-pptx-icon="check" '
        'style="position:absolute;left:104px;top:140px;width:24px;height:24px;color:#10B981"></div>'
        '</div>'
    )

    output = _normalize_slide_html(html)

    assert 'data-pptx-shape="rounded_rect"' not in output
    assert 'data-pptx-type="icon"' not in output


def test_slide_html_normalizer_keeps_table_inside_backing_card() -> None:
    html = (
        '<div data-slide="1" style="position:absolute;left:0;top:0;width:960px;height:540px">'
        '<div data-pptx-type="shape" data-pptx-shape="rounded_rect" '
        'style="position:absolute;left:80px;top:120px;width:760px;height:250px;'
        'background-color:#FFFFFF;border:1px solid #E5E7EB"></div>'
        '<div data-pptx-type="table" data-pptx-table-data=\'{"headers":["구간","역할"],"rows":[["개발","회귀"],["운영","모니터링"]]}\' '
        'style="position:absolute;left:104px;top:150px;width:712px;height:190px"></div>'
        '</div>'
    )

    elements = parse_slide_html(_normalize_slide_html(html))
    card = next(element for element in elements if element.pptx_type == "shape")
    table = next(element for element in elements if element.pptx_type == "table")

    assert table.position["left"] >= card.position["left"]
    assert table.position["top"] >= card.position["top"]
    assert table.position["left"] + table.position["width"] <= card.position["left"] + card.position["width"]
    assert table.position["top"] + table.position["height"] <= card.position["top"] + card.position["height"]


def test_slide_html_normalizer_forces_dark_text_on_light_card() -> None:
    html = (
        '<div data-slide="1" style="position:absolute;left:0;top:0;width:960px;height:540px">'
        '<div data-pptx-type="shape" data-pptx-shape="rounded_rect" '
        'style="position:absolute;left:80px;top:120px;width:360px;height:150px;'
        'background-color:#FFF7ED;border:1px solid #FDBA74"></div>'
        '<div data-pptx-type="textbox" style="position:absolute;left:112px;top:150px;'
        'width:300px;height:72px;font-size:18px;font-weight:700;color:#FFFFFF">'
        'LLM 테스트</div>'
        '</div>'
    )

    output = _normalize_slide_html(html)
    text = next(element for element in parse_slide_html(output) if element.pptx_type == "textbox")

    assert text.styles["color"].lower() not in {"#ffffff", "ffffff", "#fff"}


def test_pptx_text_parser_preserves_inline_special_character_flow() -> None:
    html = (
        '<div data-slide="1">'
        '<div data-pptx-type="textbox" style="position:absolute;left:40px;top:80px;'
        'width:820px;height:48px;font-size:28px;color:#111827">'
        'Evaluation의 3가지 층위: 개발 ▶ 스테이징 ▶ 프로덕션</div>'
        '</div>'
    )

    [element] = parse_slide_html(html)

    assert element.text_content == "Evaluation의 3가지 층위: 개발 ▶ 스테이징 ▶ 프로덕션"


def test_fallback_layout_keeps_many_unconnected_nodes_from_overlapping() -> None:
    graph = _parse_mermaid(
        "graph LR\n"
        "  A[Client]\n  B[CDN]\n  C[Load Balancer]\n  D[Web]\n"
        "  E[API]\n  F[Cache]\n  G[DB Primary]\n  H[DB Standby]\n"
        "  I[Object Storage]\n  J[Monitoring]"
    )
    positions = _layout_nodes(graph["nodes"], graph["edges"], width=1024, height=640)
    coords = list(positions.values())

    for index, (x1, y1) in enumerate(coords):
        for x2, y2 in coords[index + 1:]:
            assert abs(x1 - x2) >= 80 or abs(y1 - y2) >= 80


def test_mermaid_chain_edges_are_preserved_for_fallback_layout() -> None:
    graph = _parse_mermaid(
        "graph LR\n"
        "  U[Users] --> CDN[CDN] --> ALB[Load Balancer] --> APP[App] --> DB[(DB)]"
    )

    assert [edge["from"] for edge in graph["edges"]] == ["U", "CDN", "ALB", "APP"]
    assert [edge["to"] for edge in graph["edges"]] == ["CDN", "ALB", "APP", "DB"]


async def test_diagrams_visual_asset_renders_and_maps_to_pptx(
    tmp_path: Path,
    monkeypatch,
) -> None:
    asset = {
        "id": "aws_arch",
        "slide_index": 2,
        "method": METHOD_DIAGRAMS,
        "title": "AWS 3-Tier",
        "description": "AWS 3-tier architecture",
        "diagrams_provider": "aws",
        "diagrams_direction": "LR",
        "diagrams_clusters": [
            {"id": "region", "label": "AWS Region", "parent": ""},
            {"id": "vpc", "label": "VPC", "parent": "region"},
            {"id": "private_data", "label": "Private Data", "parent": "vpc"},
        ],
        "diagrams_nodes": [
            {"id": "U", "label": "Users", "provider": "generic", "service": "users"},
            {"id": "ALB", "label": "ALB", "provider": "aws", "service": "alb", "cluster": "vpc"},
            {
                "id": "APP",
                "label": "App Tier",
                "provider": "aws",
                "service": "ec2",
                "cluster": "vpc",
            },
            {
                "id": "DB",
                "label": "RDS",
                "provider": "aws",
                "service": "rds",
                "cluster": "private_data",
            },
        ],
        "diagrams_edges": [
            {"from": "U", "to": "ALB", "label": "HTTPS"},
            {"from": "ALB", "to": "APP"},
            {"from": "APP", "to": "DB", "label": "SQL"},
        ],
        "mermaid": (
            "graph LR\n"
            "  U[Users] --> ALB[ALB]\n"
            "  ALB --> APP[App Tier]\n"
            "  APP --> DB[(RDS)]"
        ),
        "placement": {"x": 80, "y": 100, "w": 520, "h": 320},
    }

    def fake_native_renderer(asset: dict, output_path: Path) -> dict:
        Image.new("RGB", (1800, 1100), (255, 255, 255)).save(output_path)
        return {"renderer": "diagrams", "renderer_package": "diagrams"}

    monkeypatch.setattr(
        "src.formats.pptx.agents.nodes.visual_asset_planner._render_diagrams_asset",
        fake_native_renderer,
    )

    rendered = await _render_asset(asset, tmp_path)
    assert rendered is not None
    assert Path(rendered["path"]).exists()
    assert Path(rendered["diagrams_topology_path"]).exists()
    with Image.open(rendered["path"]) as image:
        assert image.size[0] >= 1800
        assert image.size[1] >= 1100
    topology = Path(rendered["diagrams_topology_path"]).read_text(encoding="utf-8")
    assert '"provider": "aws"' in topology
    assert '"service": "rds"' in topology
    assert '"from": "APP"' in topology

    html = (
        '<div data-slide="1">'
        f'<div data-pptx-type="image" data-pptx-image-path="{rendered["path"]}" '
        'style="position:absolute;left:40px;top:100px;width:300px;height:180px"></div>'
        '</div>'
    )
    preview = _embed_local_images(html)
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    for element in parse_slide_html(html):
        CSStoOOXMLEngine()._add_element(slide, element)

    assert "data:image/png;base64" in preview
    assert slide.shapes[0].shape_type == MSO_SHAPE_TYPE.PICTURE


def test_user_attached_images_are_added_to_generation_message() -> None:
    buffer = BytesIO()
    Image.new("RGB", (4, 4), (12, 34, 56)).save(buffer, format="PNG")

    references = _build_user_reference_images(
        {
            "_image_attachments": [
                {
                    "content": buffer.getvalue(),
                    "filename": "reference.png",
                    "mime_type": "image/png",
                }
            ]
        }
    )
    human_content = _build_human_content_with_images("Generate slide", references)

    assert isinstance(human_content, list)
    assert human_content[0] == {"type": "text", "text": "Generate slide"}
    assert any(
        part.get("image_url", {}).get("url", "").startswith("data:image/png;base64,")
        for part in human_content
        if isinstance(part, dict)
    )


def test_user_attached_images_are_not_broadcast_to_content_slides() -> None:
    references = [
        {
            "label": "cover-reference.png",
            "mime_type": "image/png",
            "source": "user",
            "content": b"image",
        }
    ]

    assert _user_reference_images_for_slide(references, 1, "cover") == references
    assert _user_reference_images_for_slide(references, 2, "content") == []


def test_targeted_user_attached_images_apply_to_matching_content_slide() -> None:
    references = [
        {
            "label": "slide 3 issue.png",
            "mime_type": "image/png",
            "source": "user",
            "target_slide_index": 3,
            "content": b"image",
        }
    ]

    assert _user_reference_images_for_slide(references, 2, "content") == []
    assert _user_reference_images_for_slide(references, 3, "content") == references
